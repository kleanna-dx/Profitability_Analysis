#!/usr/bin/env python3
"""
Profitability Analysis - User Management API Manual PPT Generator
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# Paths
IMAGES_DIR = os.path.join(os.path.dirname(__file__), 'images')
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), 'Profitability_Analysis_User_API_Manual.pptx')

# Colors
BLUE = RGBColor(41, 98, 255)
DARK_BLUE = RGBColor(25, 60, 160)
LIGHT_BLUE = RGBColor(230, 240, 255)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 30, 30)
GRAY = RGBColor(100, 100, 100)
DARK_BG = RGBColor(40, 44, 52)
CODE_TEXT = RGBColor(220, 220, 220)
GREEN = RGBColor(34, 139, 34)
RED = RGBColor(220, 50, 50)
ORANGE = RGBColor(230, 150, 30)
METHOD_POST = RGBColor(73, 160, 73)
METHOD_PUT = RGBColor(230, 150, 30)
METHOD_DELETE = RGBColor(220, 60, 60)
METHOD_GET = RGBColor(41, 98, 255)
TABLE_HEADER = RGBColor(41, 98, 255)


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, font_name='Malgun Gothic', alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=12, color=BLACK, font_name='Malgun Gothic', bold=False, line_spacing=1.2):
    """Add multiline text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.3)
    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=10):
    """Add a code block with dark background"""
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()
    shape.shadow.inherit = False
    
    # Code text
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_right = Pt(12)
    tf.margin_bottom = Pt(8)
    
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = CODE_TEXT
        p.font.name = 'Consolas'
        p.space_after = Pt(2)


def add_table(slide, left, top, width, rows_data, col_widths, headers):
    """Add a styled table"""
    num_rows = len(rows_data) + 1  # +1 for header
    num_cols = len(headers)
    
    tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(0.4 * num_rows))
    tbl = tbl_shape.table
    
    # Set column widths
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    
    # Header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Malgun Gothic'
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Data rows
    for row_idx, row in enumerate(rows_data):
        for col_idx, val in enumerate(row):
            cell = tbl.cell(row_idx + 1, col_idx)
            cell.text = val
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = BLACK
            p.font.name = 'Malgun Gothic'
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return tbl_shape


def add_method_badge(slide, left, top, method, url):
    """Add HTTP method badge + URL"""
    colors = {'GET': METHOD_GET, 'POST': METHOD_POST, 'PUT': METHOD_PUT, 'DELETE': METHOD_DELETE}
    color = colors.get(method, BLUE)
    
    # Method badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.9), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = Pt(4)
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.text = method
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Consolas'
    p.alignment = PP_ALIGN.CENTER
    
    # URL text
    add_text_box(slide, left + Inches(1.0), top, Inches(7.5), Inches(0.35),
                 url, font_size=12, color=BLACK, font_name='Consolas')


def add_note_box(slide, left, top, width, text, font_size=10):
    """Add a highlighted note box"""
    height = Inches(0.6)
    # Blue left border + light bg
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    border.fill.solid()
    border.fill.fore_color.rgb = BLUE
    border.line.fill.background()
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.08), top, width - Inches(0.08), height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BLUE
    bg.line.fill.background()
    
    tf = bg.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = DARK_BLUE
    p.font.name = 'Malgun Gothic'


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ==================== SLIDE 1: COVER ====================
    slide = prs.slides.add_slide(blank_layout)
    
    # Blue gradient band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2), Inches(10), Inches(3))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    
    add_text_box(slide, Inches(0), Inches(2.5), Inches(10), Inches(0.8),
                 'User Management API Manual', font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(3.3), Inches(10), Inches(0.5),
                 'Profitability Analysis System', font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_multiline_text(slide, Inches(3), Inches(5.5), Inches(4), Inches(1.5), [
        'Version: v1.0',
        'Date: 2026-05-21',
        'Author: Kleannara DX Team',
        'For: Groupware Integration Team',
    ], font_size=11, color=GRAY)

    # ==================== SLIDE 2: TOC ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 'Table of Contents', font_size=24, bold=True, color=DARK_BLUE)
    
    toc_items = [
        '1. Overview (개요)',
        '2. Authentication (인증)',
        '3. API Endpoint Summary (API 목록)',
        '4. API 1: 전체 사용자 조회 (GET /api/users)',
        '5. API 2: 개별 사용자 조회 (GET /api/users/:userId)',
        '6. API 3: 대량 사용자 생성 (POST /api/users/bulk)',
        '7. API 4: 대량 사용자 수정 (PUT /api/users/bulk)',
        '8. API 5: 대량 사용자 비활성화 (DELETE /api/users/bulk)',
        '9. 공통 응답 코드 / 데이터 모델',
    ]
    add_multiline_text(slide, Inches(1), Inches(1.2), Inches(8), Inches(5.5),
                       toc_items, font_size=14, color=BLACK)

    # ==================== SLIDE 3: OVERVIEW ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '1. Overview (개요)', font_size=24, bold=True, color=DARK_BLUE)
    
    add_multiline_text(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(1.2), [
        'Profitability Analysis 시스템의 사용자 관리 API입니다.',
        '그룹웨어 시스템에서 사용자 정보를 동기화하기 위한 CRUD API를 제공합니다.',
    ], font_size=12, color=BLACK)
    
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(3), Inches(0.4),
                 'Base URL', font_size=13, bold=True, color=DARK_BLUE)
    add_code_block(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.5),
                   'https://3000-i4uffs666lw2u0w31ha6d-b9b802c4.sandbox.novita.ai', font_size=10)
    
    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(3), Inches(0.4),
                 'API Endpoints (5개)', font_size=13, bold=True, color=DARK_BLUE)
    add_code_block(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(1.8),
                   'GET    /api/users          # 전체 사용자 조회\n'
                   'GET    /api/users/:userId  # 개별 사용자 조회\n'
                   'POST   /api/users/bulk     # 대량 사용자 생성\n'
                   'PUT    /api/users/bulk     # 대량 사용자 수정\n'
                   'DELETE /api/users/bulk     # 대량 사용자 비활성화', font_size=11)
    
    add_note_box(slide, Inches(0.5), Inches(5.9), Inches(9),
                 'NOTE: 생성/수정/삭제는 동일 URL(/api/users/bulk)이며, HTTP Method(POST/PUT/DELETE)로 구분합니다.')

    # ==================== SLIDE 4: AUTHENTICATION ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '2. Authentication (인증)', font_size=24, bold=True, color=DARK_BLUE)
    
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.5),
                 '모든 사용자 관리 API는 아래 헤더를 반드시 포함해야 합니다.', font_size=12, color=BLACK)
    
    add_table(slide, Inches(0.5), Inches(1.8), Inches(9),
              [
                  ['Content-Type', 'application/json', '요청 데이터 형식'],
                  ['X-API-KEY', 'gw-kleannara-2026-secure-api-key', 'API 인증 키 (필수)'],
              ],
              [Inches(2.2), Inches(4.3), Inches(2.5)],
              ['Header', 'Value', 'Description'])
    
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(3), Inches(0.4),
                 '인증 실패 시 응답', font_size=13, bold=True, color=DARK_BLUE)
    add_code_block(slide, Inches(0.5), Inches(3.9), Inches(9), Inches(1.3),
                   '{\n'
                   '    "success": false,\n'
                   '    "error": "API Key가 유효하지 않습니다.",\n'
                   '    "code": "INVALID_API_KEY"\n'
                   '}', font_size=11)

    # ==================== SLIDE 5: API SUMMARY ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '3. API Endpoint Summary', font_size=24, bold=True, color=DARK_BLUE)
    
    add_table(slide, Inches(0.3), Inches(1.1), Inches(9.4),
              [
                  ['1', 'GET', '/api/users', '전체 사용자 조회 (페이지네이션, 필터)'],
                  ['2', 'GET', '/api/users/:userId', '개별 사용자 조회'],
                  ['3', 'POST', '/api/users/bulk', '대량 사용자 생성 (등록)'],
                  ['4', 'PUT', '/api/users/bulk', '대량 사용자 수정'],
                  ['5', 'DELETE', '/api/users/bulk', '대량 사용자 비활성화 (삭제)'],
              ],
              [Inches(0.5), Inches(1.2), Inches(3.5), Inches(4.2)],
              ['#', 'Method', 'URL', 'Description'])
    
    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.4),
                 'URL 구분 (동일 URL, 다른 Method)', font_size=13, bold=True, color=DARK_BLUE)
    
    add_table(slide, Inches(0.5), Inches(4.3), Inches(9),
              [
                  ['POST', '/api/users/bulk', '생성 (Create) - 새 사용자 등록'],
                  ['PUT', '/api/users/bulk', '수정 (Update) - 기존 사용자 정보 변경'],
                  ['DELETE', '/api/users/bulk', '비활성화 (Soft Delete) - is_active=0 처리'],
              ],
              [Inches(1.5), Inches(3.5), Inches(4)],
              ['HTTP Method', 'URL', 'Action'])

    # ==================== SLIDE 6: GET ALL USERS ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '4. 전체 사용자 조회', font_size=24, bold=True, color=DARK_BLUE)
    add_method_badge(slide, Inches(0.5), Inches(0.95), 'GET', '/api/users?page=1&limit=500')
    
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(0.4),
                 'Query Parameters', font_size=13, bold=True, color=DARK_BLUE)
    
    add_table(slide, Inches(0.3), Inches(1.9), Inches(9.4),
              [
                  ['page', 'number', '1', '페이지 번호 (1부터)'],
                  ['limit', 'number', '50', '페이지당 건수 (최소 1, 최대 500)'],
                  ['is_active', 'number', '-', '활성 상태 필터 (1=활성, 0=비활성)'],
                  ['role', 'string', '-', '권한 필터 (admin, user, viewer)'],
                  ['department', 'string', '-', '부서 필터 (정확 일치)'],
                  ['search', 'string', '-', '검색 (user_id, name, email 부분일치)'],
              ],
              [Inches(1.8), Inches(1.5), Inches(1.2), Inches(4.9)],
              ['Parameter', 'Type', 'Default', 'Description'])

    # ==================== SLIDE 7: GET ALL USERS - Screenshots ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '4. 전체 사용자 조회 - 요청 예시', font_size=20, bold=True, color=DARK_BLUE)
    
    # Image 1: Query Params
    img_path = os.path.join(IMAGES_DIR, 'img1_query_params.png')
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.0), width=Inches(4.5))
    
    # Image 2: GET all users response
    img_path = os.path.join(IMAGES_DIR, 'img2_get_all_users.png')
    slide.shapes.add_picture(img_path, Inches(5.0), Inches(1.0), width=Inches(4.7))

    # ==================== SLIDE 8: GET ALL USERS - Response ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '4. 전체 사용자 조회 - 응답 형식', font_size=20, bold=True, color=DARK_BLUE)
    
    add_code_block(slide, Inches(0.5), Inches(1.0), Inches(9), Inches(5.8),
                   '{\n'
                   '    "success": true,\n'
                   '    "data": [\n'
                   '        {\n'
                   '            "id": 1,\n'
                   '            "user_id": "admin",\n'
                   '            "name": "관리자",\n'
                   '            "email": null,\n'
                   '            "department": null,\n'
                   '            "phone": null,\n'
                   '            "position": null,\n'
                   '            "role": "admin",\n'
                   '            "is_active": 1,\n'
                   '            "sso_yn": 0,\n'
                   '            "created_at": "2026-05-20T07:09:57.000Z",\n'
                   '            "updated_at": "2026-05-20T08:24:32.000Z"\n'
                   '        },\n'
                   '        ...\n'
                   '    ],\n'
                   '    "pagination": { "page": 1, "limit": 500, "totalCount": 8, "totalPages": 1 }\n'
                   '}', font_size=11)

    # ==================== SLIDE 9: GET SINGLE USER ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '5. 개별 사용자 조회', font_size=24, bold=True, color=DARK_BLUE)
    add_method_badge(slide, Inches(0.5), Inches(0.95), 'GET', '/api/users/djseo')
    
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                 'Full URL: https://3000-i4uffs666lw2u0w31ha6d-b9b802c4.sandbox.novita.ai/api/users/djseo',
                 font_size=10, color=GRAY, font_name='Consolas')
    
    add_code_block(slide, Inches(0.3), Inches(2.0), Inches(4.5), Inches(4.5),
                   '// Response\n'
                   '{\n'
                   '  "success": true,\n'
                   '  "data": {\n'
                   '    "id": 3,\n'
                   '    "user_id": "djseo",\n'
                   '    "name": "서동준",\n'
                   '    "email": "djseo@kleannara.com",\n'
                   '    "department": "경영기획팀",\n'
                   '    "phone": "010-1234-5678",\n'
                   '    "position": "차장",\n'
                   '    "role": "user",\n'
                   '    "is_active": 1,\n'
                   '    "sso_yn": 1\n'
                   '  }\n'
                   '}', font_size=10)
    
    # Image 3: GET single user
    img_path = os.path.join(IMAGES_DIR, 'img3_get_single_user.png')
    slide.shapes.add_picture(img_path, Inches(5.0), Inches(2.0), width=Inches(4.7))

    # ==================== SLIDE 10: POST BULK CREATE ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '6. 대량 사용자 생성', font_size=24, bold=True, color=DARK_BLUE)
    add_method_badge(slide, Inches(0.5), Inches(0.95), 'POST', '/api/users/bulk')
    
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                 'Full URL: https://3000-i4uffs666lw2u0w31ha6d-b9b802c4.sandbox.novita.ai/api/users/bulk',
                 font_size=10, color=GRAY, font_name='Consolas')
    
    add_text_box(slide, Inches(0.5), Inches(1.9), Inches(4), Inches(0.4),
                 'Request Body Fields', font_size=13, bold=True, color=DARK_BLUE)
    
    add_table(slide, Inches(0.3), Inches(2.3), Inches(9.4),
              [
                  ['users[].userId', 'string', 'Yes', '사용자 ID (로그인 아이디)'],
                  ['users[].name', 'string', 'Yes', '사용자 이름'],
                  ['users[].password', 'string', 'Yes', '비밀번호'],
                  ['users[].email', 'string', 'No', '이메일'],
                  ['users[].department', 'string', 'No', '부서'],
                  ['users[].phone', 'string', 'No', '연락처'],
                  ['users[].position', 'string', 'No', '직급'],
                  ['users[].role', 'string', 'No', '권한 (admin/user/viewer, 기본: user)'],
                  ['users[].sso_yn', 'number', 'No', 'SSO 연동 여부 (0/1, 기본: 0)'],
              ],
              [Inches(2.5), Inches(1.3), Inches(1.0), Inches(4.6)],
              ['Field', 'Type', 'Required', 'Description'])

    # ==================== SLIDE 11: POST - Example ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '6. 대량 사용자 생성 - Request / Response', font_size=20, bold=True, color=DARK_BLUE)
    
    add_text_box(slide, Inches(0.3), Inches(0.9), Inches(4), Inches(0.3),
                 'Request Body', font_size=11, bold=True, color=DARK_BLUE)
    add_code_block(slide, Inches(0.3), Inches(1.2), Inches(4.6), Inches(4.0),
                   '{\n'
                   '  "users": [\n'
                   '    {\n'
                   '      "userId": "testuser1",\n'
                   '      "name": "테스트유저1",\n'
                   '      "password": "test1234!",\n'
                   '      "department": "개발팀",\n'
                   '      "position": "사원",\n'
                   '      "role": "user"\n'
                   '    },\n'
                   '    {\n'
                   '      "userId": "testuser2",\n'
                   '      "name": "테스트유저2",\n'
                   '      "password": "test1234!",\n'
                   '      "email": "test2@kleannara.com"\n'
                   '    }\n'
                   '  ]\n'
                   '}', font_size=9)
    
    add_text_box(slide, Inches(5.1), Inches(0.9), Inches(4), Inches(0.3),
                 'Response Body', font_size=11, bold=True, color=DARK_BLUE)
    add_code_block(slide, Inches(5.1), Inches(1.2), Inches(4.6), Inches(2.8),
                   '{\n'
                   '  "success": true,\n'
                   '  "totalCount": 2,\n'
                   '  "successCount": 2,\n'
                   '  "failCount": 0,\n'
                   '  "results": [\n'
                   '    {\n'
                   '      "userId": "testuser1",\n'
                   '      "status": "success",\n'
                   '      "message": "created"\n'
                   '    },\n'
                   '    {\n'
                   '      "userId": "testuser2",\n'
                   '      "status": "success",\n'
                   '      "message": "created"\n'
                   '    }\n'
                   '  ]\n'
                   '}', font_size=9)
    
    add_note_box(slide, Inches(0.3), Inches(5.5), Inches(9.4),
                 'NOTE: 비활성화된 동일 userId 존재 시 자동 재활성화(reactivated). 활성 중 동일 userId는 실패(already exists).')

    # ==================== SLIDE 12: PUT BULK UPDATE ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '7. 대량 사용자 수정', font_size=24, bold=True, color=DARK_BLUE)
    add_method_badge(slide, Inches(0.5), Inches(0.95), 'PUT', '/api/users/bulk')
    
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                 'Full URL: https://3000-i4uffs666lw2u0w31ha6d-b9b802c4.sandbox.novita.ai/api/users/bulk',
                 font_size=10, color=GRAY, font_name='Consolas')
    
    add_multiline_text(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(0.8), [
        '- userId: 수정 대상 식별 (필수)',
        '- 나머지 필드: 변경하고 싶은 항목만 전달 (전달하지 않은 필드는 기존 값 유지)',
    ], font_size=11, color=BLACK)
    
    add_table(slide, Inches(0.3), Inches(2.8), Inches(9.4),
              [
                  ['users[].userId', 'string', 'Yes', '수정 대상 사용자 ID'],
                  ['users[].name', 'string', 'No', '변경할 이름'],
                  ['users[].password', 'string', 'No', '변경할 비밀번호'],
                  ['users[].email', 'string', 'No', '변경할 이메일'],
                  ['users[].department', 'string', 'No', '변경할 부서'],
                  ['users[].phone', 'string', 'No', '변경할 연락처'],
                  ['users[].position', 'string', 'No', '변경할 직급'],
                  ['users[].role', 'string', 'No', '변경할 권한'],
                  ['users[].is_active', 'number', 'No', '활성 상태 (1=활성, 0=비활성)'],
              ],
              [Inches(2.5), Inches(1.3), Inches(1.0), Inches(4.6)],
              ['Field', 'Type', 'Required', 'Description'])

    # ==================== SLIDE 13: PUT - Screenshot ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '7. 대량 사용자 수정 - 요청 예시', font_size=20, bold=True, color=DARK_BLUE)
    
    # Image 4: PUT update
    img_path = os.path.join(IMAGES_DIR, 'img4_put_update.png')
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.0), width=Inches(9))

    # ==================== SLIDE 14: DELETE BULK ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '8. 대량 사용자 비활성화 (삭제)', font_size=24, bold=True, color=DARK_BLUE)
    add_method_badge(slide, Inches(0.5), Inches(0.95), 'DELETE', '/api/users/bulk')
    
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                 'Full URL: https://3000-i4uffs666lw2u0w31ha6d-b9b802c4.sandbox.novita.ai/api/users/bulk',
                 font_size=10, color=GRAY, font_name='Consolas')
    
    add_note_box(slide, Inches(0.5), Inches(1.9), Inches(9),
                 'IMPORTANT: 실제 삭제가 아닌 Soft Delete (is_active: 1→0). PUT API로 is_active: 1 설정 시 복구 가능.')
    
    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(4), Inches(0.3),
                 'Request Body', font_size=12, bold=True, color=DARK_BLUE)
    add_code_block(slide, Inches(0.5), Inches(3.1), Inches(4), Inches(0.9),
                   '{\n'
                   '    "userIds": ["testuser1", "testuser2"]\n'
                   '}', font_size=11)
    
    add_text_box(slide, Inches(5.2), Inches(2.8), Inches(4), Inches(0.3),
                 'Response Body', font_size=12, bold=True, color=DARK_BLUE)
    add_code_block(slide, Inches(5.2), Inches(3.1), Inches(4.5), Inches(2.5),
                   '{\n'
                   '  "success": true,\n'
                   '  "totalCount": 1,\n'
                   '  "successCount": 1,\n'
                   '  "failCount": 0,\n'
                   '  "results": [\n'
                   '    {\n'
                   '      "userId": "testuser1",\n'
                   '      "status": "success",\n'
                   '      "message": "deactivated"\n'
                   '    }\n'
                   '  ]\n'
                   '}', font_size=9)
    
    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(4), Inches(0.3),
                 '복구 방법 (PUT 사용)', font_size=12, bold=True, color=DARK_BLUE)
    add_code_block(slide, Inches(0.5), Inches(4.6), Inches(4), Inches(1.3),
                   'PUT /api/users/bulk\n'
                   '{\n'
                   '  "users": [\n'
                   '    { "userId": "testuser1",\n'
                   '      "is_active": 1 }\n'
                   '  ]\n'
                   '}', font_size=9)

    # ==================== SLIDE 15: DELETE - Screenshot ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '8. 대량 사용자 비활성화 - 요청 예시', font_size=20, bold=True, color=DARK_BLUE)
    
    # Image 5: DELETE
    img_path = os.path.join(IMAGES_DIR, 'img5_delete.png')
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.0), width=Inches(9))

    # ==================== SLIDE 16: RESPONSE CODES + DATA MODEL ====================
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 '9. 공통 응답 코드 / 데이터 모델', font_size=24, bold=True, color=DARK_BLUE)
    
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(4), Inches(0.3),
                 'HTTP Response Codes', font_size=12, bold=True, color=DARK_BLUE)
    add_table(slide, Inches(0.3), Inches(1.4), Inches(9.4),
              [
                  ['200', 'OK', '요청 성공'],
                  ['400', 'Bad Request', '잘못된 요청 (필수 필드 누락)'],
                  ['401', 'Unauthorized', 'API Key 인증 실패'],
                  ['404', 'Not Found', '사용자를 찾을 수 없음'],
                  ['500', 'Server Error', '서버 내부 오류'],
              ],
              [Inches(1.5), Inches(2.5), Inches(5.4)],
              ['Code', 'Status', 'Description'])
    
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(4), Inches(0.3),
                 'Users Table Schema', font_size=12, bold=True, color=DARK_BLUE)
    add_table(slide, Inches(0.3), Inches(4.4), Inches(9.4),
              [
                  ['user_id', 'VARCHAR(50)', 'Yes', '로그인 아이디 (UNIQUE)'],
                  ['name', 'VARCHAR(100)', 'Yes', '사용자 이름'],
                  ['password', 'VARCHAR(255)', 'No*', '비밀번호 (SSO는 NULL)'],
                  ['email / department / phone / position', 'VARCHAR', 'No', '부가 정보'],
                  ['role', 'ENUM', 'Yes', 'admin / user / viewer'],
                  ['is_active', 'TINYINT', 'Yes', '1=활성, 0=비활성'],
                  ['sso_yn', 'TINYINT', 'Yes', '1=SSO, 0=일반'],
              ],
              [Inches(3.5), Inches(2.0), Inches(1.0), Inches(2.9)],
              ['Column', 'Type', 'Required', 'Description'])

    # ==================== SAVE ====================
    prs.save(OUTPUT_PATH)
    print(f'PPT generated: {OUTPUT_PATH}')
    print(f'File size: {os.path.getsize(OUTPUT_PATH) / 1024:.1f} KB')


if __name__ == '__main__':
    build_ppt()
